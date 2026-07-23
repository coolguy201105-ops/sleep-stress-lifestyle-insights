"""
Generate a color-themed presentation deck summarizing the project.

Run:
    python make_slides.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Project_Slides.pptx"

# ---------------------------------------------------------------------------
# Theme (matches the Streamlit app's violet palette)
# ---------------------------------------------------------------------------
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_DARK = RGBColor(0x5B, 0x21, 0xB6)
VIOLET_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
GRAY_TEXT = RGBColor(0x1F, 0x29, 0x37)
GRAY_MUTED = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=GRAY_TEXT,
             bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri",
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=GRAY_TEXT,
                 bullet_color=VIOLET, font="Calibri", space_after=10, bold_lead=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = f"\u25CF  {lead}"
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = bullet_color
            r1.font.name = font
            if rest:
                r2 = p.add_run()
                r2.text = f"  {rest}"
                r2.font.size = Pt(size)
                r2.font.color.rgb = color
                r2.font.name = font
        else:
            r1 = p.add_run()
            r1.text = f"\u25CF  {item}"
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = font
    return box


def slide_header(slide, title, subtitle=None, band_color=VIOLET):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), band_color)
    add_text(slide, Inches(0.55), Inches(0.12), Inches(12), Inches(0.7),
              title, size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.25), Inches(12), Inches(0.4),
                  subtitle, size=14, color=GRAY_MUTED, italic=True)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None,
              header_color=VIOLET, font_size=13):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = w
    for c, h in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = VIOLET_LIGHT if r % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(font_size - 1)
                    run.font.color.rgb = GRAY_TEXT
    return gtable


def add_stat_card(slide, left, top, width, height, value, label, color=VIOLET):
    card = add_rect(slide, left, top, width, height, VIOLET_LIGHT)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    add_text(slide, left, top + Inches(0.12), width, Inches(0.65), value,
              size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + height - Inches(0.55), width - Inches(0.2),
              Inches(0.5), label, size=11, color=GRAY_MUTED, align=PP_ALIGN.CENTER)


def add_footer(slide, page_text):
    add_text(slide, Inches(0.55), SLIDE_H - Inches(0.42), Inches(8), Inches(0.35),
              "Sleep, Stress & Lifestyle Insights", size=10, color=GRAY_MUTED)
    add_text(slide, SLIDE_W - Inches(2.0), SLIDE_H - Inches(0.42), Inches(1.5), Inches(0.35),
              page_text, size=10, color=GRAY_MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, VIOLET)
add_rect(s, 0, Inches(4.5), SLIDE_W, Inches(0.06), WHITE)
add_text(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.3),
          "Sleep, Stress & Lifestyle Insights", size=44, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.35), Inches(11.3), Inches(0.7),
          "A Data-Driven Biomedical Research Project for the Irvine Community",
          size=20, color=VIOLET_LIGHT, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1), Inches(4.85), Inches(11.3), Inches(0.5),
          "Baseline machine learning models \u2022 Kaggle sleep & lifestyle datasets \u2022 Interactive web app",
          size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
          "github.com/coolguy201105-ops/sleep-stress-lifestyle-insights",
          size=12, color=VIOLET_LIGHT, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 2: Problem & Motivation
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Problem & Community Relevance")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(7.4), Inches(4.8), [
    ("Poor sleep and high stress", "are widespread concerns among students and families."),
    ("Consequences are real:", "reduced academic performance, weakened immune function, and long-term health risk."),
    ("No local Irvine dataset exists", "but national/international research shows the same risk factors likely apply here."),
    ("This project's goal:", "find which lifestyle and physiological factors most strongly relate to poor sleep and stress, and share the findings through an interactive tool."),
], size=17)
card = add_rect(s, Inches(8.3), Inches(1.7), Inches(4.4), Inches(4.6), VIOLET_LIGHT)
add_text(s, Inches(8.55), Inches(1.95), Inches(3.9), Inches(0.5), "Research Question",
          size=16, bold=True, color=VIOLET)
add_text(s, Inches(8.55), Inches(2.5), Inches(3.9), Inches(3.6),
          "How are sleep duration, sleep quality, and lifestyle factors (activity, "
          "caffeine, alcohol, exercise, BMI) associated with stress, sleep disorders, "
          "and sleep efficiency?",
          size=15, color=GRAY_TEXT, line_spacing=1.2)
add_footer(s, "2")

# ---------------------------------------------------------------------------
# Slide 3: Data Sources
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Data Sources", "Five public Kaggle datasets reviewed \u2014 two used for modeling")
headers = ["Dataset", "Size", "Used?", "Provides"]
rows = [
    ["Sleep Health & Lifestyle", "374 \u00d7 13", "\u2705 Main", "Sleep, stress, activity, BMI, vitals, disorder diagnosis"],
    ["Sleep Efficiency", "452 \u00d7 15", "\u2705 Supporting", "Sleep efficiency, sleep stages, awakenings, caffeine/alcohol"],
    ["Animal sleep dataset", "62 \u00d7 8", "\u274c Excluded", "Non-human species data \u2014 not relevant"],
    ["Wearable tech sleep quality", "1,000 \u00d7 9", "\u274c Not used", "Simulated sensor data, no demographics"],
    ["Health Sleep Statistics", "100 \u00d7 12", "\u274c Not used", "Small sample, overlaps main dataset, no stress column"],
]
add_table(s, Inches(0.55), Inches(1.55), Inches(12.2), Inches(3.6), headers, rows,
          col_widths=[Inches(3.0), Inches(1.5), Inches(1.6), Inches(6.1)], font_size=13)
add_text(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.3),
          "The two datasets used were analyzed separately \u2014 never merged row-by-row \u2014 "
          "since they were collected independently with different columns and sample sizes.",
          size=14, italic=True, color=GRAY_MUTED)
add_footer(s, "3")

# ---------------------------------------------------------------------------
# Slide 4: Methods
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Methods", "A 3-tier modeling approach for every target variable")
steps = [
    ("1. Data cleaning", "Remove duplicates, fix a hidden data-loading bug, convert timestamps to numeric minutes."),
    ("2. Exploratory analysis", "Visualize relationships \u2014 e.g., stress by sleep-duration group."),
    ("3. Predictive modeling", "Train & compare 3 models per target on an identical 75/25 train-test split."),
    ("4. Interpretation", "Rank Random Forest feature importances to find the strongest predictors."),
    ("5. Deployment", "Package everything into a public, interactive Streamlit web app."),
]
top = Inches(1.65)
for i, (title, desc) in enumerate(steps):
    y = top + Inches(1.05) * i
    add_rect(s, Inches(0.6), y, Inches(0.5), Inches(0.5), VIOLET)
    add_text(s, Inches(0.6), y, Inches(0.5), Inches(0.5), str(i + 1), size=20, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.3), y - Inches(0.05), Inches(3.0), Inches(0.5), title, size=16,
              bold=True, color=VIOLET_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.5), y - Inches(0.05), Inches(8.2), Inches(0.6), desc, size=14,
              color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, "4")

# ---------------------------------------------------------------------------
# Slide 5: The Three Models
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Three Models, Compared Head-to-Head")
cards = [
    ("Dummy Baseline", "Ignores every input. Predicts the average (or most common class).\nPurpose: a sanity-check floor \u2014 nothing to beat, everything to compare against.", GRAY_MUTED),
    ("Linear / Logistic Regression", "A weighted sum of all features.\nPurpose: the interpretable, explainable model \u2014 each feature gets one clear coefficient.", BLUE),
    ("Random Forest", "Hundreds of decision trees, averaged together.\nPurpose: the best-effort model \u2014 captures non-linear effects and feature interactions.", VIOLET),
]
w = Inches(3.95)
for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.55) + (w + Inches(0.2)) * i
    add_rect(s, x, Inches(1.6), w, Inches(0.65), color)
    add_text(s, x, Inches(1.6), w, Inches(0.65), title, size=16, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box = add_rect(s, x, Inches(2.35), w, Inches(3.6), VIOLET_LIGHT)
    add_text(s, x + Inches(0.2), Inches(2.55), w - Inches(0.4), Inches(3.3), desc, size=14,
              color=GRAY_TEXT, line_spacing=1.25)
add_footer(s, "5")

# ---------------------------------------------------------------------------
# Slide 6: Results Overview
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Results Overview", "Random Forest wins on every target \u2014 but by how much varies")
headers = ["Target", "Best Model", "Best Score", "vs. Dummy Baseline"]
rows = [
    ["Stress Level (0\u201310)", "Random Forest", "R\u00b2 = 0.94", "Error cut ~84%"],
    ["Quality of Sleep (1\u201310)", "Random Forest", "R\u00b2 = 0.94", "Error cut ~89%"],
    ["Sleep Disorder (3-class)", "Random Forest", "66.7% accuracy", "+12 pts accuracy"],
    ["Sleep Efficiency (0\u20131)", "Random Forest", "R\u00b2 = 0.86", "Error cut ~68%"],
    ["Awakenings (count/night)", "Random Forest", "R\u00b2 = 0.58", "Error cut ~39%"],
]
add_table(s, Inches(0.55), Inches(1.6), Inches(12.2), Inches(3.2), headers, rows,
          col_widths=[Inches(3.6), Inches(2.8), Inches(2.8), Inches(3.0)], font_size=14)
add_stat_card(s, Inches(0.9), Inches(5.1), Inches(2.7), Inches(1.5), "0.94", "Best R\u00b2\n(Stress & Sleep Quality)", GREEN)
add_stat_card(s, Inches(3.9), Inches(5.1), Inches(2.7), Inches(1.5), "66.7%", "Best classification\naccuracy", AMBER)
add_stat_card(s, Inches(6.9), Inches(5.1), Inches(2.7), Inches(1.5), "5 of 5", "Targets beat the\nbaseline", BLUE)
add_stat_card(s, Inches(9.9), Inches(5.1), Inches(2.7), Inches(1.5), "132\u2013452", "Rows analyzed\nper dataset", VIOLET)
add_footer(s, "6")

# ---------------------------------------------------------------------------
# Slide 7: Deep Dive - Stress & Sleep Quality
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Deep Dive: Stress & Sleep Quality", band_color=GREEN)
headers = ["Model", "Stress \u2014 MAE / R\u00b2", "Sleep Quality \u2014 MAE / R\u00b2"]
rows = [
    ["Dummy", "1.41 / ~0.00", "1.04 / ~0.00"],
    ["Linear Regression", "0.30 / 0.71", "0.17 / 0.84"],
    ["Random Forest", "0.22 / 0.94", "0.12 / 0.94"],
]
add_table(s, Inches(0.55), Inches(1.6), Inches(7.4), Inches(2.4), headers, rows,
          col_widths=[Inches(2.2), Inches(2.6), Inches(2.6)], header_color=GREEN, font_size=14)
add_text(s, Inches(0.55), Inches(4.3), Inches(7.4), Inches(2.5),
          "Both targets are strongly predictable \u2014 Random Forest's clear edge over Linear "
          "Regression suggests non-linear or interaction effects (e.g., stress may not scale "
          "evenly with sleep duration alone).",
          size=15, color=GRAY_TEXT, line_spacing=1.3)
box = add_rect(s, Inches(8.3), Inches(1.6), Inches(4.4), Inches(5.2), VIOLET_LIGHT)
add_text(s, Inches(8.55), Inches(1.8), Inches(3.9), Inches(0.5), "Top Predictors", size=16, bold=True, color=GREEN)
add_bullets(s, Inches(8.55), Inches(2.35), Inches(3.9), Inches(4.3), [
    ("Stress level:", "sleep quality, heart rate, sleep duration"),
    ("Sleep quality:", "sleep duration (dominant), stress level, heart rate"),
], size=14, bullet_color=GREEN)
add_footer(s, "7")

# ---------------------------------------------------------------------------
# Slide 8: Deep Dive - Sleep Disorder Classification
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Deep Dive: Sleep Disorder Classification", band_color=AMBER)
headers = ["Model", "Accuracy", "Balanced Accuracy", "Macro F1"]
rows = [
    ["Dummy (\"always None\")", "54.5%", "33.3%", "0.24"],
    ["Logistic Regression", "63.6%", "60.4%", "0.59"],
    ["Random Forest", "66.7%", "62.2%", "0.61"],
]
add_table(s, Inches(0.55), Inches(1.6), Inches(8.0), Inches(2.4), headers, rows,
          col_widths=[Inches(2.8), Inches(1.7), Inches(2.0), Inches(1.5)], header_color=AMBER, font_size=14)
add_text(s, Inches(0.55), Inches(4.3), Inches(7.9), Inches(2.6),
          "A modest but real result: Random Forest beats the dummy baseline by ~12 points of "
          "accuracy and nearly 30 points of balanced accuracy. Still, 63\u201367% on a 3-class "
          "problem is an early signal, not a diagnostic tool.",
          size=15, color=GRAY_TEXT, line_spacing=1.3)
box = add_rect(s, Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.2), VIOLET_LIGHT)
box.line.color.rgb = RED
box.line.width = Pt(1.5)
add_text(s, Inches(9.1), Inches(1.8), Inches(3.4), Inches(0.9), "Data-Cleaning Fix", size=15, bold=True, color=RED)
add_text(s, Inches(9.1), Inches(2.4), Inches(3.4), Inches(4.2),
          "An earlier version misread the text \"None\" as a missing value, silently dropping "
          "219 of 374 rows and inflating classification accuracy. Fixed by disabling pandas' "
          "default NA behavior \u2014 numbers shown here are corrected.",
          size=13, color=GRAY_TEXT, line_spacing=1.25)
add_footer(s, "8")

# ---------------------------------------------------------------------------
# Slide 9: Deep Dive - Sleep Efficiency & Awakenings
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Deep Dive: Sleep Efficiency & Awakenings", band_color=BLUE)
headers = ["Model", "Efficiency \u2014 MAE / R\u00b2", "Awakenings \u2014 MAE / R\u00b2"]
rows = [
    ["Dummy", "0.120 / ~0.00", "1.20 / ~0.00"],
    ["Linear Regression", "0.050 / 0.80", "0.80 / 0.46"],
    ["Random Forest", "0.038 / 0.86", "0.73 / 0.58"],
]
add_table(s, Inches(0.55), Inches(1.6), Inches(7.4), Inches(2.4), headers, rows,
          col_widths=[Inches(2.2), Inches(2.6), Inches(2.6)], header_color=BLUE, font_size=14)
add_text(s, Inches(0.55), Inches(4.3), Inches(7.4), Inches(2.6),
          "Awakenings is the hardest target (R\u00b2 0.58 at best) \u2014 likely because important "
          "drivers like noise, room temperature, or undiagnosed conditions aren't captured "
          "in this dataset.",
          size=15, color=GRAY_TEXT, line_spacing=1.3)
box = add_rect(s, Inches(8.3), Inches(1.6), Inches(4.4), Inches(5.2), VIOLET_LIGHT)
add_text(s, Inches(8.55), Inches(1.8), Inches(3.9), Inches(0.5), "Interesting Finding", size=16, bold=True, color=BLUE)
add_text(s, Inches(8.55), Inches(2.35), Inches(3.9), Inches(4.3),
          "Sleep-structure metrics (light/deep sleep %, prior awakenings) drove these "
          "predictions more than lifestyle habits. Caffeine, alcohol, and smoking mattered "
          "less than expected in this sample.",
          size=14, color=GRAY_TEXT, line_spacing=1.25)
add_footer(s, "9")

# ---------------------------------------------------------------------------
# Slide 10: Community Takeaways
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Community Takeaways", band_color=VIOLET)
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.0), Inches(5.0), [
    ("Sleep duration & quality are the biggest levers", "for lower stress in this data \u2014 more consistent, higher-quality sleep is linked to meaningfully lower stress."),
    ("Sleep disorder risk is hard to flag from lifestyle data alone", "(mid-60% accuracy) \u2014 a promising early signal, not a substitute for clinical screening."),
    ("Sleep efficiency problems may need different interventions", "than simple lifestyle changes, since caffeine/alcohol/exercise had smaller effects than expected."),
    ("All results are associations, not causation", "\u2014 this is observational data, and the app/tool is educational only, not medical advice."),
], size=18, space_after=18)
add_footer(s, "10")

# ---------------------------------------------------------------------------
# Slide 11: The Web App
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Try It Yourself: The Interactive Web App", band_color=GREEN)
add_bullets(s, Inches(0.6), Inches(1.7), Inches(7.2), Inches(4.8), [
    ("Overview", "project intro and dataset summary"),
    ("Explore the Data", "interactive charts for both datasets"),
    ("Model Performance", "live comparisons of all 3 model types for all 5 targets"),
    ("Try the Predictor", "enter your own stats, get real-time predictions with plain-language feedback"),
    ("Full Report", "the complete methodology and results, in one place"),
], size=17, bullet_color=GREEN)
box = add_rect(s, Inches(8.3), Inches(1.7), Inches(4.4), Inches(4.6), GREEN)
add_text(s, Inches(8.55), Inches(2.0), Inches(3.9), Inches(0.5), "Live App", size=16, bold=True, color=WHITE)
add_text(s, Inches(8.55), Inches(2.6), Inches(3.9), Inches(1.2),
          "sleep-stress-lifestyle-insights\n.streamlit.app",
          size=16, bold=True, color=WHITE, line_spacing=1.2)
add_text(s, Inches(8.55), Inches(4.0), Inches(3.9), Inches(2.0),
          "(replace with your exact deployed link if different)",
          size=12, italic=True, color=VIOLET_LIGHT)
add_footer(s, "11")

# ---------------------------------------------------------------------------
# Slide 12: Limitations & Ethics
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Limitations & Ethics", band_color=RED)
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.0), Inches(5.0), [
    ("Not local data", "datasets are public and not collected in Irvine \u2014 findings represent general patterns, not verified local statistics."),
    ("Association, not causation", "all datasets are observational; results describe correlations, not proof of cause and effect."),
    ("Smaller effective sample", "the main lifestyle dataset shrank from 374 to 132 unique rows after removing exact duplicates."),
    ("Educational tool only", "this project and its predictor are not a diagnostic or medical device, and do not replace professional medical advice."),
], size=18, bullet_color=RED, space_after=18)
add_footer(s, "12")

# ---------------------------------------------------------------------------
# Slide 13: Conclusion & Next Steps
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s)
slide_header(s, "Conclusion & Next Steps", band_color=VIOLET)
add_text(s, Inches(0.6), Inches(1.6), Inches(12.0), Inches(1.1),
          "Everyday, self-reported lifestyle and physiological factors \u2014 especially sleep "
          "duration and quality \u2014 carry a strong, measurable relationship with stress and "
          "sleep health outcomes in this data.",
          size=18, color=GRAY_TEXT, line_spacing=1.3)
add_bullets(s, Inches(0.6), Inches(3.0), Inches(12.0), Inches(3.5), [
    ("Validate", "these patterns against a larger or local sample if one becomes available."),
    ("Extend", "the wearable-sensor dataset into modeling once paired with real demographic data."),
    ("Translate", "the strongest, most actionable findings into a one-page community handout for Irvine students and families."),
], size=17, space_after=16)
add_footer(s, "13")

# ---------------------------------------------------------------------------
# Slide 14: Thank You
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, VIOLET)
add_text(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1.0),
          "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.9), Inches(11.3), Inches(0.6),
          "Questions & Discussion", size=20, color=VIOLET_LIGHT, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1), Inches(5.4), Inches(11.3), Inches(0.5),
          "github.com/coolguy201105-ops/sleep-stress-lifestyle-insights",
          size=13, color=VIOLET_LIGHT, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
